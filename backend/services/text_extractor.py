"""
Text Extractor — Shared text extraction library for files and HTML strings.

Pure functions — no database, no MemoryStore, no Chalie services.
Used by the document upload pipeline and the `read` innate skill
(URL fetch + local file read).

TEXT ONLY. An image carries no text to extract — it is DESCRIBED by
``services.image_description.ImageDescription`` (the vision tool's describe core),
and the document pipeline routes it there on mime before ever reaching this module.
Images are rejected loudly here rather than silently plain-read (see extract_text).

Supported formats (heavy-library imports are lazy):
    PDF(pdfplumber), DOCX(python-docx), PPTX(python-pptx),
    HTML(trafilatura), plain text/markdown.
"""

import logging
import mimetypes
from typing import TYPE_CHECKING

if TYPE_CHECKING:
    from collections.abc import Sequence
    from typing import Protocol

    from docx.document import Document as _DocxDocument
    from pdfplumber.page import Page as _PdfPage

    # python-pptx is Any under mypy (no py.typed), so a minimal structural
    # type stands in for the slide object the helper actually uses.
    class _PptxParagraph(Protocol):
        text: str

    class _PptxTextFrame(Protocol):
        paragraphs: "Sequence[_PptxParagraph]"

    class _PptxShape(Protocol):
        has_text_frame: bool
        text_frame: "_PptxTextFrame"

    class _PptxSlide(Protocol):
        shapes: "Sequence[_PptxShape]"

logger = logging.getLogger(__name__)


# ─── Public API ──────────────────────────────────────────────────────────────

def extract_text(file_path: str, mime_type: str | None = None) -> str:
    """Dispatches to format-specific extractors by MIME type.

    Raises ValueError for any image/* type: an image has no text to extract, and
    the plain-read fallback below would otherwise decode its binary as UTF-8
    mojibake and return that as a 'successful' extraction. Callers that may hold
    an image must fork on mime first and route it to ImageDescription().
    """
    if not mime_type:
        mime_type = detect_mime_type(file_path)

    if mime_type.startswith('image/'):
        raise ValueError(
            f"extract_text cannot read '{mime_type}': images are described, not extracted. "
            "Route images to services.image_description.ImageDescription()."
        )

    extractors = {
        'application/pdf': _extract_pdf,
        'application/vnd.openxmlformats-officedocument.wordprocessingml.document': _extract_docx,
        'application/vnd.openxmlformats-officedocument.presentationml.presentation': _extract_pptx,
        'text/html': _extract_html_file,
        'text/plain': _extract_plain,
        'text/markdown': _extract_plain,
    }

    extractor = extractors.get(mime_type)
    if extractor:
        return extractor(file_path)

    # Fallback: any text/* type (code files, CSV, etc.)
    if mime_type.startswith('text/'):
        return _extract_plain(file_path)

    logger.warning(f"[TEXT EXTRACTOR] Unsupported mime type '{mime_type}' — attempting plain read")
    return _extract_plain(file_path)


def extract_html(html: str, url: str | None = None) -> str:
    """
    Extract clean, readable text from an HTML string via trafilatura.

    Returns empty string if extraction fails or produces no content.
    """
    if not html or not html.strip():
        return ''

    try:
        import trafilatura
        content = trafilatura.extract(html, url=url, include_comments=False, include_links=True)
        if content and content.strip():
            return content.strip()
    except Exception as e:
        logger.debug(f'[TEXT EXTRACTOR] trafilatura failed: {e}')

    return ''



def detect_mime_type(file_path: str) -> str:
    """
    Detect MIME type from file extension using stdlib mimetypes.

    Falls back to 'text/plain' for unknown extensions.
    """
    mime_type, _ = mimetypes.guess_type(file_path)
    return mime_type or 'text/plain'


# ─── Format-specific extractors (internal) ───────────────────────────────────

def _pdf_page_text(page: "_PdfPage") -> str:
    """Extract one PDF page's text, appending any detected tables as
    pipe-joined rows. Extracted from _extract_pdf."""
    page_text = page.extract_text() or ''

    tables = page.extract_tables()
    if tables:
        for table in tables:
            rows = []
            for row in table:
                cells = [str(cell or '').strip() for cell in row]
                rows.append(' | '.join(cells))
            page_text += '\n' + '\n'.join(rows)

    return page_text


def _extract_pdf(path: str) -> str:
    """Extract text from PDF using pdfplumber with table detection."""
    try:
        import pdfplumber
    except ImportError:
        logger.error('[TEXT EXTRACTOR] pdfplumber not installed — cannot extract PDF')
        return ''

    try:
        pages = []
        with pdfplumber.open(path) as pdf:
            for i, page in enumerate(pdf.pages):
                page_text = _pdf_page_text(page)

                if page_text.strip():
                    pages.append(f"[Page {i + 1}]\n{page_text.strip()}")

        return '\n\n'.join(pages)

    except Exception as e:
        logger.exception(f'[TEXT EXTRACTOR] PDF extraction failed: {e}')
        return ''


def _docx_paragraph_text(doc: "_DocxDocument", element: object) -> str | None:
    """Return the formatted text for the paragraph matching ``element``
    (heading-prefixed with markdown '#'s if the paragraph uses a Heading
    style), or ``None`` if empty or no matching paragraph is found.
    Extracted from _extract_docx."""
    for para in doc.paragraphs:
        if para._element == element:
            text = para.text.strip()
            if not text:
                return None
            if para.style and para.style.name.startswith('Heading'):
                level_str = para.style.name.replace('Heading ', '').replace('Heading', '1')
                try:
                    level = int(level_str)
                except ValueError:
                    level = 1
                return f"{'#' * level} {text}"
            return text
    return None


def _docx_table_text(doc: "_DocxDocument", element: object) -> str | None:
    """Return the pipe-joined row text for the table matching ``element``,
    or ``None`` if no matching table is found. Extracted from
    _extract_docx."""
    for table in doc.tables:
        if table._element == element:
            rows = []
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                rows.append(' | '.join(cells))
            return '\n'.join(rows)
    return None


def _extract_docx(path: str) -> str:
    """Extract text from DOCX with paragraph and table support."""
    try:
        from docx import Document
    except ImportError:
        logger.error('[TEXT EXTRACTOR] python-docx not installed — cannot extract DOCX')
        return ''

    try:
        doc = Document(path)
        parts = []

        for element in doc.element.body:
            tag = element.tag.split('}')[-1] if '}' in element.tag else element.tag
            if tag == 'p':
                text = _docx_paragraph_text(doc, element)
                if text is not None:
                    parts.append(text)
            elif tag == 'tbl':
                table_text = _docx_table_text(doc, element)
                if table_text is not None:
                    parts.append(table_text)

        return '\n\n'.join(parts)

    except Exception as e:
        logger.exception(f'[TEXT EXTRACTOR] DOCX extraction failed: {e}')
        return ''


def _pptx_slide_texts(slide: "_PptxSlide") -> list[str]:
    """Collect all non-empty paragraph texts across a slide's text-frame
    shapes. Extracted from _extract_pptx."""
    texts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if text:
                    texts.append(text)
    return texts


def _extract_pptx(path: str) -> str:
    """Extract text from PowerPoint slides as labelled sections."""
    try:
        from pptx import Presentation
    except ImportError:
        logger.error('[TEXT EXTRACTOR] python-pptx not installed — cannot extract PPTX')
        return ''

    try:
        prs = Presentation(path)
        slides = []

        for i, slide in enumerate(prs.slides):
            texts = _pptx_slide_texts(slide)
            if texts:
                slides.append(f"[Slide {i + 1}]\n" + '\n'.join(texts))

        return '\n\n'.join(slides)

    except Exception as e:
        logger.exception(f'[TEXT EXTRACTOR] PPTX extraction failed: {e}')
        return ''


def _extract_html_file(path: str) -> str:
    """Extract text from an HTML file on disk. Delegates to extract_html()."""
    try:
        with open(path, 'r', encoding='utf-8', errors='replace') as f:
            raw = f.read()
        return extract_html(raw)
    except Exception as e:
        logger.exception(f'[TEXT EXTRACTOR] HTML file extraction failed: {e}')
        return ''


def _extract_plain(path: str) -> str:
    """Read a plain text file with UTF-8 encoding, replacing undecodable bytes."""
    try:
        with open(path, 'r', encoding='utf-8', errors='replace') as f:
            return f.read()
    except Exception as e:
        logger.exception(f'[TEXT EXTRACTOR] Plain text read failed: {e}')
        return ''
